#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import os
import glob
import re
import tempfile # For temporary files
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import mermaid # For generating SVG from Mermaid code
import cairosvg # For converting SVG to PNG

def create_single_slide_presentation(md_content, output_path):
    """
    Parses markdown content to create a single-slide presentation.
    The slide will contain the Gantt chart as an image at the top and all other
    relevant information below, with a font size of 12.
    """
    lines = md_content.split('\n')
    
    title = "Work Status Report" # Default title
    mermaid_code = ""
    other_content_lines = []

    # --- 1. Parse Markdown Content ---

    # Find the main title (H1)
    for line in lines:
        if line.startswith('# '):
            title = line[2:].strip()
            break
            
    # Extract the content of the Mermaid block
    in_mermaid = False
    mermaid_block_start_line = -1
    mermaid_block_end_line = -1
    for i, line in enumerate(lines):
        if line.strip() == '```mermaid':
            in_mermaid = True
            mermaid_block_start_line = i
            continue
        if in_mermaid and line.strip() == '```':
            in_mermaid = False
            mermaid_block_end_line = i
            break
        if in_mermaid:
            mermaid_code += line + '\n'
    
    # Gather all other relevant content lines
    for i, line in enumerate(lines):
        # Skip title, comments, horizontal rules, and the mermaid block itself
        if line.startswith('# '): continue
        if line.strip().startswith('<!--'): continue
        if line.strip() == '---': continue
        if mermaid_block_start_line <= i <= mermaid_block_end_line: continue
        
        # Keep the line if it has content
        if line.strip():
            other_content_lines.append(line)

    # --- 2. Create the Presentation ---
    
    prs = Presentation()
    # Use "Title and Content" layout
    slide_layout = prs.slide_layouts[1] 
    slide = prs.slides.add_slide(slide_layout)
    
    # Set the slide title
    slide.shapes.title.text = title
    # Apply font to title
    slide.shapes.title.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
    
    # Get the body text frame
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()  # Clear default text
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    # --- Handle Mermaid Gantt Chart ---
    gantt_image_path = None
    if mermaid_code.strip():
        print("Attempting to render Mermaid Gantt chart to image...")
        temp_svg_path = None
        temp_png_path = None
        try:
            # Create temporary file paths
            with tempfile.NamedTemporaryFile(delete=False, suffix=".svg") as temp_svg_file:
                temp_svg_path = temp_svg_file.name
            with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as temp_png_file:
                temp_png_path = temp_png_file.name

            # Inject font family into Mermaid code
            # This ensures Mermaid.js uses a Chinese-compatible font during SVG generation
            mermaid_code_with_font = f"%%{{init: {{'theme': 'default', 'fontFamily': 'Microsoft YaHei, sans-serif'}} }}%%\n{mermaid_code}"

            # 1. Generate SVG from Mermaid code directly to file
            mermaid.Mermaid(mermaid_code_with_font).to_svg(path=temp_svg_path)
            
            # 2. Convert SVG to PNG
            cairosvg.svg2png(url=temp_svg_path, write_to=temp_png_path)
            gantt_image_path = temp_png_path
            print(f"Mermaid Gantt chart rendered to: {gantt_image_path}")

        except Exception as e:
            print(f"Error rendering Mermaid chart to image: {e}")
            print("Falling back to text version of Gantt chart.")
            # Fallback to text if image generation fails
            p_mermaid = tf.add_paragraph()
            run = p_mermaid.add_run()
            run.text = mermaid_code
            run.font.size = Pt(12)
            run.font.name = 'Microsoft YaHei' # Use a Chinese font for the chart text
        finally:
            if temp_svg_path and os.path.exists(temp_svg_path):
                os.remove(temp_svg_path)
            # temp_png_path will be cleaned up after prs.save()

    if gantt_image_path:
        # Position the image at the top of the slide
        left = Inches(0.5)
        top = Inches(1.5) # Below the title
        width = Inches(9)
        height = Inches(3) # Adjust as needed

        pic = slide.shapes.add_picture(gantt_image_path, left, top, width=width, height=height)
        
        # Adjust the text frame position to be below the image
        tf.top = top + height + Inches(0.2) # 0.2 inch padding
        tf.height = prs.slide_height - tf.top - Inches(0.5) # Remaining height

    # Add a separator line if an image was inserted, or if falling back to text
    if gantt_image_path or mermaid_code.strip():
        p_sep = tf.add_paragraph()
        run = p_sep.add_run()
        run.text = '\n' + ('-' * 40) + '\n'
        run.font.size = Pt(12)
        run.font.name = 'Microsoft YaHei' # Apply font to separator

    # Add all other content
    for line_text in other_content_lines:
        p = tf.add_paragraph()
        
        cleaned_text = line_text.strip() # Remove all leading/trailing whitespace

        # Remove Links from Sub-items: finds [text](url) and replaces it with just text
        cleaned_text = re.sub(r'\[(.*?)\]\(.*\)', r'\1', cleaned_text)

        run = p.add_run() # Add a run to the paragraph
        run.font.size = Pt(12)
        run.font.name = 'Microsoft YaHei'
        
        # Bold and Level 1 for Status Headings
        if cleaned_text.endswith('Projects'): # Heuristic to identify status headings
            run.text = cleaned_text
            run.font.bold = True
            p.level = 0 # Highest level
        elif cleaned_text.startswith('* '):
            run.text = cleaned_text.lstrip('* ').strip() # Remove bullet and any extra spaces
            p.level = 2
        elif cleaned_text.startswith('- '):
            run.text = cleaned_text.lstrip('- ').strip() # Remove bullet and any extra spaces
            p.level = 1
        else:
            run.text = cleaned_text
            p.level = 0
        
    print(f"Saving single-slide presentation to: {output_path}")
    prs.save(output_path)

    # Clean up the generated PNG file after saving the presentation
    if gantt_image_path and os.path.exists(gantt_image_path):
        os.remove(gantt_image_path)


def main():
    """
    Main function to find the markdown file, parse it, and generate a presentation.
    """
    base_path = os.path.dirname(os.path.abspath(__file__))
    people_path = os.path.join(base_path, 'people')

    try:
        md_file = glob.glob(os.path.join(people_path, '*.md'))[0]
    except IndexError:
        print("Error: No markdown file found in the /people directory.")
        return

    print(f"Processing file: {md_file}")

    with open(md_file, 'r', encoding='utf-8') as f:
        md_content = f.read()
    
    output_filename = os.path.splitext(os.path.basename(md_file))[0] + '.pptx'
    output_path = os.path.join(people_path, output_filename)
    
    create_single_slide_presentation(md_content, output_path)
    print("Presentation created successfully.")

if __name__ == '__main__':
    main()